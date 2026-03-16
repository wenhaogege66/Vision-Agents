"""PPTConvertService 单元测试。"""

import io
from unittest.mock import MagicMock, patch

import pytest
from fastapi import HTTPException
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.services.ppt_convert_service import PPTConvertService


# ── Helpers ──────────────────────────────────────────────────


def _mock_supabase() -> MagicMock:
    """Create a mock Supabase client."""
    return MagicMock()


def _make_pptx_bytes(num_slides: int = 3, add_text: bool = True) -> bytes:
    """Create a valid .pptx file in memory with the given number of slides.

    Args:
        num_slides: Number of slides to create.
        add_text: Whether to add text content to each slide.

    Returns:
        Bytes of a valid .pptx file.
    """
    prs = Presentation()
    for i in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        if add_text:
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            )
            tf = txBox.text_frame
            tf.text = f"Slide {i + 1} Title"
            p = tf.add_paragraph()
            p.text = f"Content for slide {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_empty_pptx_bytes() -> bytes:
    """Create a valid .pptx file with zero slides."""
    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _setup_storage_download(sb: MagicMock, content: bytes):
    """Mock sb.storage.from_(bucket).download(path) to return content."""
    bucket = MagicMock()
    bucket.download.return_value = content
    sb.storage.from_.return_value = bucket
    return bucket


def _setup_storage_download_and_upload(sb: MagicMock, content: bytes):
    """Mock storage for both download and upload."""
    bucket = MagicMock()
    bucket.download.return_value = content
    bucket.upload.return_value = None
    sb.storage.from_.return_value = bucket
    return bucket


def _setup_table_chain(sb: MagicMock, data=None):
    """Set up a chainable mock for sb.table(name).update()...execute()."""
    chain = MagicMock()
    chain.execute.return_value = MagicMock(data=data or [])
    for method in ("update", "eq"):
        getattr(chain, method).return_value = chain
    sb.table.return_value = chain
    return chain


# ── convert_to_images Tests ──────────────────────────────────


class TestConvertToImages:
    """Tests for PPTConvertService.convert_to_images."""

    @pytest.mark.asyncio
    async def test_returns_correct_number_of_image_paths(self):
        """Converted image count should equal the number of slides."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=5)
        _setup_storage_download_and_upload(sb, pptx_bytes)

        svc = PPTConvertService(sb)
        paths = await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")

        assert len(paths) == 5

    @pytest.mark.asyncio
    async def test_image_paths_follow_expected_format(self):
        """Each image path should follow {project_id}/{material_type}/images/page_NNN.png."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=3)
        _setup_storage_download_and_upload(sb, pptx_bytes)

        svc = PPTConvertService(sb)
        paths = await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")

        assert paths == [
            "proj-1/text_ppt/images/page_001.png",
            "proj-1/text_ppt/images/page_002.png",
            "proj-1/text_ppt/images/page_003.png",
        ]

    @pytest.mark.asyncio
    async def test_single_slide_ppt(self):
        """A PPT with one slide should produce exactly one image."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=1)
        _setup_storage_download_and_upload(sb, pptx_bytes)

        svc = PPTConvertService(sb)
        paths = await svc.convert_to_images("proj-1/presentation_ppt/v2_deck.pptx")

        assert len(paths) == 1
        assert paths[0] == "proj-1/presentation_ppt/images/page_001.png"

    @pytest.mark.asyncio
    async def test_empty_ppt_raises_400(self):
        """A PPT with zero slides should raise HTTPException 400."""
        sb = _mock_supabase()
        pptx_bytes = _make_empty_pptx_bytes()
        _setup_storage_download_and_upload(sb, pptx_bytes)

        svc = PPTConvertService(sb)
        with pytest.raises(HTTPException) as exc_info:
            await svc.convert_to_images("proj-1/text_ppt/v1_empty.pptx")
        assert exc_info.value.status_code == 400
        assert "不包含任何幻灯片" in str(exc_info.value.detail)

    @pytest.mark.asyncio
    async def test_invalid_pptx_raises_500(self):
        """Non-PPTX content should raise HTTPException 500."""
        sb = _mock_supabase()
        _setup_storage_download_and_upload(sb, b"not-a-pptx-file")

        svc = PPTConvertService(sb)
        with pytest.raises(HTTPException) as exc_info:
            await svc.convert_to_images("proj-1/text_ppt/v1_bad.pptx")
        assert exc_info.value.status_code == 500
        assert "解析失败" in str(exc_info.value.detail)

    @pytest.mark.asyncio
    async def test_download_failure_raises_500(self):
        """Storage download failure should raise HTTPException 500."""
        sb = _mock_supabase()
        bucket = MagicMock()
        bucket.download.side_effect = Exception("network error")
        sb.storage.from_.return_value = bucket

        svc = PPTConvertService(sb)
        with pytest.raises(HTTPException) as exc_info:
            await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")
        assert exc_info.value.status_code == 500
        assert "下载" in str(exc_info.value.detail)

    @pytest.mark.asyncio
    async def test_upload_failure_raises_500(self):
        """Storage upload failure should raise HTTPException 500."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=1)
        bucket = MagicMock()
        bucket.download.return_value = pptx_bytes
        bucket.upload.side_effect = Exception("storage full")
        sb.storage.from_.return_value = bucket

        svc = PPTConvertService(sb)
        with pytest.raises(HTTPException) as exc_info:
            await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")
        assert exc_info.value.status_code == 500
        assert "上传图像失败" in str(exc_info.value.detail)

    @pytest.mark.asyncio
    async def test_uploaded_images_are_valid_png(self):
        """Each uploaded image should be a valid PNG."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=2)
        bucket = MagicMock()
        bucket.download.return_value = pptx_bytes
        uploaded_images: list[bytes] = []

        def capture_upload(path, file, file_options=None):
            uploaded_images.append(file)

        bucket.upload.side_effect = capture_upload
        sb.storage.from_.return_value = bucket

        svc = PPTConvertService(sb)
        await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")

        assert len(uploaded_images) == 2
        for img_bytes in uploaded_images:
            # Verify it's a valid PNG by opening with Pillow
            img = Image.open(io.BytesIO(img_bytes))
            assert img.format == "PNG"
            assert img.width > 0
            assert img.height > 0

    @pytest.mark.asyncio
    async def test_slides_without_text_still_produce_images(self):
        """Slides with no text should still generate valid images."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=2, add_text=False)
        _setup_storage_download_and_upload(sb, pptx_bytes)

        svc = PPTConvertService(sb)
        paths = await svc.convert_to_images("proj-1/text_ppt/v1_blank.pptx")

        assert len(paths) == 2

    @pytest.mark.asyncio
    async def test_storage_upload_called_with_correct_paths(self):
        """Verify each upload call uses the correct storage path."""
        sb = _mock_supabase()
        pptx_bytes = _make_pptx_bytes(num_slides=2)
        bucket = MagicMock()
        bucket.download.return_value = pptx_bytes
        sb.storage.from_.return_value = bucket

        svc = PPTConvertService(sb)
        await svc.convert_to_images("proj-1/text_ppt/v1_slides.pptx")

        upload_calls = bucket.upload.call_args_list
        assert len(upload_calls) == 2

        paths_used = []
        for call in upload_calls:
            path = call.kwargs.get("path") or call[1].get("path") or call[0][0]
            paths_used.append(path)

        assert "proj-1/text_ppt/images/page_001.png" in paths_used
        assert "proj-1/text_ppt/images/page_002.png" in paths_used


# ── update_material_image_paths Tests ────────────────────────


class TestUpdateMaterialImagePaths:
    """Tests for PPTConvertService.update_material_image_paths."""

    @pytest.mark.asyncio
    async def test_updates_image_paths_in_db(self):
        """Should call table update with correct image_paths JSON."""
        sb = _mock_supabase()
        chain = _setup_table_chain(sb)

        svc = PPTConvertService(sb)
        paths = ["proj-1/text_ppt/images/page_001.png", "proj-1/text_ppt/images/page_002.png"]
        await svc.update_material_image_paths("mat-123", paths)

        sb.table.assert_called_once_with("project_materials")
        chain.update.assert_called_once()
        chain.eq.assert_called_once_with("id", "mat-123")

    @pytest.mark.asyncio
    async def test_db_failure_raises_500(self):
        """Database update failure should raise HTTPException 500."""
        sb = _mock_supabase()
        chain = MagicMock()
        chain.update.return_value = chain
        chain.eq.return_value = chain
        chain.execute.side_effect = Exception("db error")
        sb.table.return_value = chain

        svc = PPTConvertService(sb)
        with pytest.raises(HTTPException) as exc_info:
            await svc.update_material_image_paths("mat-123", [])
        assert exc_info.value.status_code == 500
        assert "更新图像路径失败" in str(exc_info.value.detail)


# ── _build_images_base_path Tests ────────────────────────────


class TestBuildImagesBasePath:
    """Tests for the static path builder."""

    def test_standard_path(self):
        result = PPTConvertService._build_images_base_path(
            "proj-1/text_ppt/v1_slides.pptx"
        )
        assert result == "proj-1/text_ppt/images"

    def test_presentation_ppt_path(self):
        result = PPTConvertService._build_images_base_path(
            "abc-123/presentation_ppt/v3_deck.pptx"
        )
        assert result == "abc-123/presentation_ppt/images"

    def test_invalid_path_raises_500(self):
        with pytest.raises(HTTPException) as exc_info:
            PPTConvertService._build_images_base_path("no_slash")
        assert exc_info.value.status_code == 500
