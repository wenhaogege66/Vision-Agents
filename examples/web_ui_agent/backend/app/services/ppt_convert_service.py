"""PPT转图像服务：将PPT文件每页转换为PNG图像并上传到Supabase Storage。

使用 python-pptx 提取幻灯片内容，Pillow 生成带有幻灯片文本的代表性图像。
生产环境可替换为 LibreOffice headless 方案以获得更高保真度的渲染效果。
"""

import io
import json
import logging
import tempfile
from pathlib import Path

from fastapi import HTTPException
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from supabase import Client

logger = logging.getLogger(__name__)

# Supabase Storage bucket 名称（与 material_service 保持一致）
STORAGE_BUCKET = "materials"

# 默认图像尺寸（像素），对应标准16:9幻灯片
DEFAULT_IMAGE_WIDTH = 1280
DEFAULT_IMAGE_HEIGHT = 720

# 图像样式
BG_COLOR = (255, 255, 255)
TEXT_COLOR = (51, 51, 51)
TITLE_COLOR = (33, 33, 33)
SLIDE_NUM_COLOR = (153, 153, 153)
BORDER_COLOR = (220, 220, 220)


class PPTConvertService:
    """将PPT文件转换为页面图像并上传到Supabase Storage。"""

    def __init__(self, supabase: Client) -> None:
        self._sb = supabase

    async def convert_to_images(self, file_path: str) -> list[str]:
        """将PPT文件转换为页面图像列表。

        流程：
        1. 从 Supabase Storage 下载PPT文件
        2. 使用 python-pptx 解析幻灯片
        3. 为每页生成PNG图像（提取文本内容渲染到图像上）
        4. 上传每张PNG到 Supabase Storage
        5. 返回图像路径列表

        Args:
            file_path: Supabase Storage 中的PPT文件路径，
                       格式如 "{project_id}/{material_type}/v{version}_{filename}"

        Returns:
            图像路径列表，如 ["{project_id}/{material_type}/images/page_001.png", ...]

        Raises:
            HTTPException: 下载、解析或上传失败时抛出
        """
        # 1. 下载PPT文件
        ppt_bytes = self._download_file(file_path)

        # 2. 解析PPT
        try:
            prs = Presentation(io.BytesIO(ppt_bytes))
        except Exception as exc:
            logger.exception("解析PPT文件失败: %s", file_path)
            raise HTTPException(
                status_code=500, detail=f"PPT文件解析失败: {exc}"
            ) from exc

        if len(prs.slides) == 0:
            raise HTTPException(status_code=400, detail="PPT文件不包含任何幻灯片")

        # 计算图像存储的基础路径
        images_base = self._build_images_base_path(file_path)

        # 3 & 4. 为每页生成图像并上传
        image_paths: list[str] = []
        slide_width = prs.slide_width or Emu(12192000)  # 默认10英寸
        slide_height = prs.slide_height or Emu(6858000)  # 默认7.5英寸

        for idx, slide in enumerate(prs.slides, start=1):
            # 提取幻灯片文本内容
            texts = self._extract_slide_texts(slide)

            # 生成图像
            img_bytes = self._render_slide_image(
                slide_number=idx,
                total_slides=len(prs.slides),
                texts=texts,
                slide_width=slide_width,
                slide_height=slide_height,
            )

            # 上传图像
            page_name = f"page_{idx:03d}.png"
            storage_path = f"{images_base}/{page_name}"
            self._upload_image(storage_path, img_bytes)
            image_paths.append(storage_path)

        logger.info(
            "PPT转换完成: %s -> %d 张图像", file_path, len(image_paths)
        )
        return image_paths

    async def update_material_image_paths(
        self, material_id: str, image_paths: list[str]
    ) -> None:
        """更新 project_materials 记录的 image_paths 字段。

        Args:
            material_id: 材料记录ID
            image_paths: 图像路径列表
        """
        try:
            self._sb.table("project_materials").update(
                {"image_paths": json.dumps(image_paths)}
            ).eq("id", material_id).execute()
        except Exception as exc:
            logger.exception("更新材料图像路径失败: %s", material_id)
            raise HTTPException(
                status_code=500, detail=f"更新图像路径失败: {exc}"
            ) from exc

    # ── 内部方法 ──────────────────────────────────────────────

    def _download_file(self, file_path: str) -> bytes:
        """从 Supabase Storage 下载文件。"""
        try:
            response = self._sb.storage.from_(STORAGE_BUCKET).download(file_path)
            return response
        except Exception as exc:
            logger.exception("从 Storage 下载文件失败: %s", file_path)
            raise HTTPException(
                status_code=500, detail=f"下载PPT文件失败: {exc}"
            ) from exc

    @staticmethod
    def _build_images_base_path(file_path: str) -> str:
        """根据PPT文件路径构建图像存储基础路径。

        输入: "{project_id}/{material_type}/v1_slides.pptx"
        输出: "{project_id}/{material_type}/images"
        """
        parts = file_path.split("/")
        if len(parts) < 2:
            raise HTTPException(
                status_code=500,
                detail=f"无效的文件路径格式: {file_path}",
            )
        # 取 project_id 和 material_type（前两段）
        return f"{parts[0]}/{parts[1]}/images"

    @staticmethod
    def _extract_slide_texts(slide) -> list[str]:
        """从幻灯片中提取所有文本内容。"""
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        return texts

    @staticmethod
    def _render_slide_image(
        slide_number: int,
        total_slides: int,
        texts: list[str],
        slide_width,
        slide_height,
    ) -> bytes:
        """使用 Pillow 渲染幻灯片为PNG图像。

        创建白色背景图像，绘制幻灯片编号和提取的文本内容。
        """
        # 计算宽高比并确定图像尺寸
        if slide_width and slide_height:
            aspect = int(slide_width) / int(slide_height)
            img_w = DEFAULT_IMAGE_WIDTH
            img_h = int(img_w / aspect)
        else:
            img_w = DEFAULT_IMAGE_WIDTH
            img_h = DEFAULT_IMAGE_HEIGHT

        img = Image.new("RGB", (img_w, img_h), BG_COLOR)
        draw = ImageDraw.Draw(img)

        # 绘制边框
        draw.rectangle(
            [0, 0, img_w - 1, img_h - 1], outline=BORDER_COLOR, width=2
        )

        # 使用默认字体（跨平台兼容）
        try:
            title_font = ImageFont.truetype("arial.ttf", 28)
            body_font = ImageFont.truetype("arial.ttf", 18)
            num_font = ImageFont.truetype("arial.ttf", 14)
        except (OSError, IOError):
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            num_font = ImageFont.load_default()

        # 绘制页码
        page_text = f"Slide {slide_number} / {total_slides}"
        draw.text((img_w - 160, img_h - 30), page_text, fill=SLIDE_NUM_COLOR, font=num_font)

        # 绘制文本内容
        y_offset = 40
        max_text_width = img_w - 80  # 左右各留40px边距

        if not texts:
            draw.text(
                (40, y_offset),
                f"[Slide {slide_number}]",
                fill=TITLE_COLOR,
                font=title_font,
            )
        else:
            # 第一行作为标题
            title = texts[0]
            if len(title) > 60:
                title = title[:57] + "..."
            draw.text((40, y_offset), title, fill=TITLE_COLOR, font=title_font)
            y_offset += 50

            # 分隔线
            draw.line([(40, y_offset), (img_w - 40, y_offset)], fill=BORDER_COLOR, width=1)
            y_offset += 20

            # 其余文本作为正文
            for text in texts[1:]:
                if y_offset > img_h - 60:
                    draw.text((40, y_offset), "...", fill=TEXT_COLOR, font=body_font)
                    break
                # 截断过长文本
                if len(text) > 80:
                    text = text[:77] + "..."
                draw.text((40, y_offset), text, fill=TEXT_COLOR, font=body_font)
                y_offset += 28

        # 导出为PNG字节
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    def _upload_image(self, storage_path: str, img_bytes: bytes) -> None:
        """上传图像到 Supabase Storage。"""
        try:
            self._sb.storage.from_(STORAGE_BUCKET).upload(
                path=storage_path,
                file=img_bytes,
                file_options={"content-type": "image/png"},
            )
        except Exception as exc:
            logger.exception("上传图像失败: %s", storage_path)
            raise HTTPException(
                status_code=500, detail=f"上传图像失败: {exc}"
            ) from exc
